from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BLUE = RGBColor(30, 58, 138)
INDIGO = RGBColor(99, 102, 241)
TEXT = RGBColor(55, 65, 81)
MUTED = RGBColor(107, 114, 128)
LIGHT = RGBColor(240, 244, 255)
WHITE = RGBColor(255, 255, 255)

DIM_LABEL = {
    "features": "核心功能",
    "pricing": "定价模式",
    "target_users": "目标用户",
    "positioning": "产品定位",
    "swot": "优劣势",
    "strategic_implications": "战略启示",
}
KEY_DIMS = ["features", "pricing", "target_users", "positioning"]


def _load_json(path: Path, default: Any) -> Any:
    if not path.exists():
        return default
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return default


def _safe_text(value: Any, fallback: str = "—") -> str:
    text = str(value or "").strip()
    return text or fallback


def _truncate(value: Any, limit: int) -> str:
    text = _safe_text(value, "")
    return text[:limit] + ("…" if len(text) > limit else "")


def _support_score(claim: dict[str, Any]) -> float:
    try:
        return float(claim.get("support_score") or 0.0)
    except (TypeError, ValueError):
        return 0.0


def _accepted_claims(claims: list[dict[str, Any]]) -> list[dict[str, Any]]:
    accepted = [claim for claim in claims if claim.get("accepted") is not False]
    return sorted(accepted, key=_support_score, reverse=True)


def _competitors(summary: dict[str, Any], claims: list[dict[str, Any]]) -> list[str]:
    raw = summary.get("competitors")
    if isinstance(raw, list):
        competitors = [_safe_text(item, "") for item in raw if _safe_text(item, "")]
    else:
        competitors = []
    if competitors:
        return competitors
    names = {
        _safe_text(claim.get("competitor_name"), "")
        for claim in claims
        if claim.get("competitor_name")
    }
    return sorted(name for name in names if name)


def _build_matrix(
    competitors: list[str],
    accepted: list[dict[str, Any]],
) -> dict[str, dict[str, str]]:
    by_comp_dim: dict[str, dict[str, str]] = {comp: {} for comp in competitors}
    for claim in accepted:
        comp = _safe_text(claim.get("competitor_name"), "")
        dim = _safe_text(claim.get("dimension"), "")
        statement = _safe_text(claim.get("statement"), "")
        if comp in by_comp_dim and dim and statement and dim not in by_comp_dim[comp]:
            by_comp_dim[comp][dim] = _truncate(statement, 80)
    return by_comp_dim


def _qa_pass_rate(claims: list[dict[str, Any]], accepted: list[dict[str, Any]]) -> str:
    if not claims:
        return "—"
    return f"{round((len(accepted) / len(claims)) * 100):.0f}%"


def _top_source_ids(evidence: list[dict[str, Any]]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for item in evidence:
        source_id = _safe_text(item.get("source_id"), "")
        if not source_id or source_id in seen:
            continue
        seen.add(source_id)
        result.append(source_id)
        if len(result) == 3:
            break
    return result


def _load_ppt_data(run_dir: Path) -> dict[str, Any]:
    run_dir = Path(run_dir)
    summary = _load_json(run_dir / "run_summary.json", {})
    claims_payload = _load_json(run_dir / "claims.json", [])
    evidence_payload = _load_json(run_dir / "evidence.json", [])
    source_summary = _load_json(run_dir / "source_summary.json", {})

    summary = summary if isinstance(summary, dict) else {}
    claims = [item for item in claims_payload if isinstance(item, dict)]
    evidence = [item for item in evidence_payload if isinstance(item, dict)]
    source_summary = source_summary if isinstance(source_summary, dict) else {}
    accepted = _accepted_claims(claims)
    competitors = _competitors(summary, accepted or claims)

    return {
        "summary": summary,
        "claims": claims,
        "accepted": accepted,
        "top_claims": accepted[:5],
        "evidence": evidence,
        "source_summary": source_summary,
        "competitors": competitors,
        "matrix": _build_matrix(competitors, accepted),
        "qa_pass_rate": _qa_pass_rate(claims, accepted),
        "top_sources": _top_source_ids(evidence),
    }


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN | None = None,
):
    shape = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    _add_text(slide, title, 0.55, 0.35, 8.9, 0.45, size=26, bold=True, color=BLUE)
    if subtitle:
        _add_text(slide, subtitle, 0.58, 0.82, 8.4, 0.28, size=10, color=MUTED)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.55),
        Inches(1.16),
        Inches(9.2),
        Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = INDIGO
    line.line.fill.background()


def _add_metric_card(slide, title: str, value: str, left: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(1.65),
        Inches(2.65),
        Inches(1.25),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = RGBColor(219, 234, 254)
    _add_text(slide, value, left + 0.18, 1.84, 2.25, 0.38, size=24, bold=True, color=BLUE)
    _add_text(slide, title, left + 0.18, 2.28, 2.25, 0.25, size=11, color=MUTED)


def _slide_cover(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(1.18),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    _add_text(slide, "竞品分析报告", 0.62, 0.28, 5.8, 0.48, size=28, bold=True, color=WHITE)

    summary = payload["summary"]
    query = _safe_text(summary.get("query"), "竞品分析")
    _add_text(slide, _truncate(query, 72), 0.72, 1.75, 9.4, 1.1, size=30, bold=True, color=BLUE)

    for index, name in enumerate(payload["competitors"][:5]):
        left = 0.78 + index * 1.65
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(3.12),
            Inches(1.42),
            Inches(0.42),
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor(229, 231, 235)
        chip.line.fill.background()
        _add_text(slide, _truncate(name, 12), left + 0.1, 3.22, 1.22, 0.18, size=10, bold=True)

    completed = _safe_text(summary.get("completed_at"), "")[:10] or "—"
    task_id = _safe_text(summary.get("task_id"), _safe_text(summary.get("run_id"), "—"))
    footer = f"{completed} | Run ID: {task_id} | 由 cs-mvp 系统生成"
    _add_text(slide, footer, 0.75, 6.78, 8.5, 0.26, size=10, color=MUTED)


def _slide_top_findings(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_title(slide, "核心结论", "按 support_score 选取前 5 条 accepted claim")
    top_claims = payload["top_claims"]
    if not top_claims:
        _add_text(slide, "暂无 accepted claims。", 0.75, 1.65, 9.0, 0.35, size=15)
        return
    for index, claim in enumerate(top_claims, start=1):
        comp = _safe_text(claim.get("competitor_name"), "跨竞品")
        dim = DIM_LABEL.get(_safe_text(claim.get("dimension"), ""), _safe_text(claim.get("dimension"), "维度"))
        statement = _truncate(claim.get("statement"), 118)
        y = 1.45 + (index - 1) * 0.88
        _add_text(slide, f"{index}. [{comp} · {dim}] {statement}", 0.78, y, 8.85, 0.54, size=14)


def _slide_matrix(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_title(slide, "横向对比矩阵", "行=维度，列=竞品，单元格来自 claims.json")
    competitors = payload["competitors"] or ["—"]
    matrix = payload["matrix"]
    rows = len(KEY_DIMS) + 1
    cols = len(competitors) + 1
    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.55),
        Inches(1.42),
        Inches(12.2),
        Inches(4.75),
    )
    table = table_shape.table
    table.columns[0].width = Inches(1.35)
    remaining = 10.85 / max(1, len(competitors))
    for col in range(1, cols):
        table.columns[col].width = Inches(remaining)

    table.cell(0, 0).text = ""
    for col, name in enumerate(competitors, start=1):
        table.cell(0, col).text = _truncate(name, 18)
        table.cell(0, col).fill.solid()
        table.cell(0, col).fill.fore_color.rgb = BLUE
    for row, dim in enumerate(KEY_DIMS, start=1):
        table.cell(row, 0).text = DIM_LABEL.get(dim, dim)
        table.cell(row, 0).fill.solid()
        table.cell(row, 0).fill.fore_color.rgb = LIGHT
        for col, comp in enumerate(competitors, start=1):
            table.cell(row, col).text = matrix.get(comp, {}).get(dim, "—")

    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Microsoft YaHei"
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = TEXT
    for col in range(1, cols):
        cell = table.rows[0].cells[col]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE


def _slide_evidence_summary(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_title(slide, "证据摘要", "来自 evidence.json 和 source_summary.json")
    source_summary = payload["source_summary"]
    valid_sources = source_summary.get("valid") or source_summary.get("fetched") or 0
    _add_metric_card(slide, "有效来源", str(valid_sources), 0.65)
    _add_metric_card(slide, "证据条数", str(len(payload["evidence"])), 3.55)
    _add_metric_card(slide, "QA 通过率", payload["qa_pass_rate"], 6.45)

    _add_text(slide, "Top source ids / URLs", 0.75, 3.45, 5.2, 0.32, size=16, bold=True, color=BLUE)
    top_sources = payload["top_sources"] or ["—"]
    for index, source in enumerate(top_sources, start=1):
        _add_text(slide, f"{index}. {_truncate(source, 95)}", 0.9, 3.9 + index * 0.42, 8.7, 0.26, size=11)


def _slide_system_note(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_title(slide, "分析说明", "PPT 由结构化 JSON artifact 自动生成")
    summary = payload["summary"]
    duration = summary.get("duration_seconds")
    duration_text = f"{float(duration):.0f}s" if isinstance(duration, int | float) else "—"
    task_id = _safe_text(summary.get("task_id"), _safe_text(summary.get("run_id"), "—"))
    lines = [
        "本报告由 cs-mvp Agent 系统自动生成。",
        "PPT 数据只读取 run 目录中的结构化 JSON artifact，不解析 report.md 或 report.html。",
        "核心结论来自 accepted claims，证据摘要来自 evidence.json，来源统计来自 source_summary.json。",
        "PPT 生成失败不会影响主流程完成状态；原始审计报告仍以 report.md / report.html 为准。",
        f"Run ID: {task_id}",
        f"分析耗时: {duration_text}",
        "总成本: —",
    ]
    for index, line in enumerate(lines):
        _add_text(slide, f"• {line}", 0.8, 1.55 + index * 0.55, 9.3, 0.34, size=14)


def build_ppt(run_dir: Path) -> Path:
    """Read JSON artifacts from run_dir and generate report.pptx."""
    run_dir = Path(run_dir)
    payload = _load_ppt_data(run_dir)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs, payload)
    _slide_top_findings(prs, payload)
    _slide_matrix(prs, payload)
    _slide_evidence_summary(prs, payload)
    _slide_system_note(prs, payload)

    pptx_path = run_dir / "report.pptx"
    prs.save(pptx_path)
    return pptx_path
