from __future__ import annotations

from pathlib import Path
from typing import Any

from jinja2 import Template

from cs_mvp.agents.ppt_builder import _load_ppt_data


SLIDE_WIDTH_PX = 1280
SLIDE_HEIGHT_PX = 720
SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5

TEMPLATE_DIR = Path(__file__).parent / "ppt_slides"

DIM_LABEL = {
    "features": "核心功能",
    "pricing": "定价模式",
    "target_users": "目标用户",
    "positioning": "产品定位",
    "swot": "优劣势",
    "strategic_implications": "战略启示",
}
KEY_DIMS = ["features", "pricing", "target_users", "positioning"]


def _safe_text(value: Any, fallback: str = "—") -> str:
    text = str(value or "").strip()
    return text or fallback


def _truncate(value: Any, limit: int) -> str:
    text = _safe_text(value, "")
    return text[:limit] + ("…" if len(text) > limit else "")


def _render_template(template_name: str, **values: Any) -> str:
    template_path = TEMPLATE_DIR / template_name
    template = Template(template_path.read_text(encoding="utf-8"))
    return template.render(**values)


def _completed_at(summary: dict[str, Any]) -> str:
    value = _safe_text(summary.get("completed_at"), "")
    return value[:19] if value else "—"


def _task_id(summary: dict[str, Any], run_dir: Path) -> str:
    return _safe_text(
        summary.get("task_id"),
        _safe_text(summary.get("run_id"), run_dir.name),
    )


def _duration_text(summary: dict[str, Any]) -> str:
    value = summary.get("duration_seconds")
    try:
        seconds = float(value)
    except (TypeError, ValueError):
        return "—"
    return f"{seconds:.0f}s"


def _render_cover(payload: dict[str, Any], run_dir: Path) -> str:
    summary = payload["summary"]
    competitors = [_truncate(item, 22) for item in payload["competitors"][:6]]
    return _render_template(
        "slide_cover.html",
        query=_truncate(summary.get("query"), 92) or "竞品分析",
        competitors=competitors,
        completed_at=_completed_at(summary),
        task_id=_task_id(summary, run_dir),
    )


def _render_findings(payload: dict[str, Any]) -> str:
    findings = []
    for index, claim in enumerate(payload["top_claims"][:5], start=1):
        dim = _safe_text(claim.get("dimension"), "")
        findings.append(
            {
                "index": index,
                "competitor": _truncate(claim.get("competitor_name"), 24),
                "dimension": DIM_LABEL.get(dim, dim or "维度"),
                "statement": _truncate(claim.get("statement"), 118),
            }
        )
    return _render_template("slide_findings.html", findings=findings)


def _render_matrix(payload: dict[str, Any]) -> str:
    competitors = payload["competitors"][:4] or ["—"]
    matrix = payload["matrix"]
    rows = []
    for dim in KEY_DIMS:
        rows.append(
            {
                "dimension": DIM_LABEL.get(dim, dim),
                "cells": [
                    _truncate(matrix.get(comp, {}).get(dim), 90)
                    for comp in competitors
                ],
            }
        )
    return _render_template(
        "slide_matrix.html",
        competitors=[_truncate(item, 18) for item in competitors],
        rows=rows,
    )


def _valid_source_count(source_summary: dict[str, Any]) -> Any:
    for key in ("valid", "fetched", "total"):
        value = source_summary.get(key)
        if value is not None:
            return value
    return 0


def _render_evidence(payload: dict[str, Any]) -> str:
    top_sources = [_truncate(source, 86) for source in payload["top_sources"][:3]]
    if not top_sources:
        top_sources = ["—"]
    return _render_template(
        "slide_evidence.html",
        valid_sources=_valid_source_count(payload["source_summary"]),
        evidence_count=len(payload["evidence"]),
        qa_pass_rate=payload["qa_pass_rate"],
        top_sources=top_sources,
    )


def _render_note(payload: dict[str, Any], run_dir: Path) -> str:
    summary = payload["summary"]
    notes = [
        "PPT 仅读取 claims.json、run_summary.json、evidence.json、source_summary.json 等结构化 artifact。",
        "核心结论来自 accepted claims，证据摘要来自 evidence.json，来源统计来自 source_summary.json。",
        "视觉版生成失败时，系统会降级到原文字版 report.pptx，不影响主流程完成。",
        "审计型报告仍以 report.md / report.html 为准，PPT 用于演示与快速汇报。",
    ]
    return _render_template(
        "slide_note.html",
        notes=notes,
        task_id=_task_id(summary, run_dir),
        duration=_duration_text(summary),
        completed_at=_completed_at(summary),
    )


def _slide_defs(payload: dict[str, Any], run_dir: Path) -> list[tuple[str, str]]:
    return [
        ("slide_cover.html", _render_cover(payload, run_dir)),
        ("slide_findings.html", _render_findings(payload)),
        ("slide_matrix.html", _render_matrix(payload)),
        ("slide_evidence.html", _render_evidence(payload)),
        ("slide_note.html", _render_note(payload, run_dir)),
    ]


def _render_slides_to_png(
    slide_defs: list[tuple[str, str]],
    run_dir: Path,
) -> list[Path]:
    """Render each HTML slide into a 1280x720 PNG with Playwright."""
    from pathlib import Path as _Path

    from playwright.sync_api import sync_playwright

    from cs_mvp.tools.fetch import _PLAYWRIGHT_EXE

    slides_dir = Path(run_dir) / "_slides"
    slides_dir.mkdir(parents=True, exist_ok=True)

    html_png_pairs: list[tuple[Path, Path]] = []
    for index, (template_name, html) in enumerate(slide_defs, start=1):
        stem = _Path(template_name).stem
        html_path = slides_dir / f"{index:02d}_{stem}.html"
        png_path = slides_dir / f"{index:02d}_{stem}.png"
        html_path.write_text(html, encoding="utf-8")
        html_png_pairs.append((html_path, png_path))

    launch_kwargs: dict[str, Any] = {"headless": True}
    exe_path = Path(_PLAYWRIGHT_EXE)
    if exe_path.exists():
        launch_kwargs["executable_path"] = str(exe_path)

    with sync_playwright() as p:
        browser = p.chromium.launch(**launch_kwargs)
        try:
            page = browser.new_page(
                viewport={"width": SLIDE_WIDTH_PX, "height": SLIDE_HEIGHT_PX},
                device_scale_factor=1,
            )
            for html_path, png_path in html_png_pairs:
                page.goto(html_path.resolve().as_uri(), wait_until="load", timeout=10_000)
                page.wait_for_timeout(500)
                page.screenshot(
                    path=str(png_path),
                    clip={
                        "x": 0,
                        "y": 0,
                        "width": SLIDE_WIDTH_PX,
                        "height": SLIDE_HEIGHT_PX,
                    },
                )
        finally:
            browser.close()

    return [png_path for _, png_path in html_png_pairs]


def _assemble_pptx(png_paths: list[Path], run_dir: Path) -> Path:
    """Assemble full-page PNG screenshots into report.pptx."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]
    for png_path in png_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png_path),
            Inches(0),
            Inches(0),
            Inches(SLIDE_WIDTH_IN),
            Inches(SLIDE_HEIGHT_IN),
        )

    pptx_path = Path(run_dir) / "report.pptx"
    prs.save(pptx_path)
    return pptx_path


def build_ppt_visual(run_dir: Path) -> Path:
    """Generate a high-fidelity PPT from HTML/CSS slide screenshots.

    The caller owns fallback behavior. If Playwright or PPT assembly fails, this
    function raises so `node_finalize` can fall back to the text PPT builder.
    """
    run_dir = Path(run_dir)
    payload = _load_ppt_data(run_dir)
    png_paths = _render_slides_to_png(_slide_defs(payload, run_dir), run_dir)
    return _assemble_pptx(png_paths, run_dir)
