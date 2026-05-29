from __future__ import annotations

import base64
import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from cs_mvp import db
from cs_mvp.agents import ppt_visual_builder
from cs_mvp.agents.ppt_visual_builder import build_ppt_visual
from cs_mvp.models import AgentRun, AnalysisTask, CompetitorInput, GraphState


_PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
)


def _write_json(path: Path, payload: object) -> None:
    path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")


@pytest.fixture()
def minimal_run_dir(tmp_path: Path) -> Path:
    _write_json(
        tmp_path / "run_summary.json",
        {
            "task_id": "T-visual-001",
            "run_id": "RUN-visual-001",
            "query": "AI code editor competitive analysis",
            "competitors": ["Cursor", "Copilot", "Tabnine"],
            "status": "completed",
            "completed_at": "2026-05-25 12:00:00",
            "duration_seconds": 300.0,
        },
    )
    _write_json(
        tmp_path / "claims.json",
        [
            {
                "claim_id": "C-001",
                "competitor_name": "Cursor",
                "dimension": "pricing",
                "statement": "Cursor offers a Pro plan for product teams.",
                "support_score": 0.85,
                "accepted": True,
                "evidence_ids": ["E-001"],
            },
            {
                "claim_id": "C-002",
                "competitor_name": "Copilot",
                "dimension": "features",
                "statement": "GitHub Copilot is deeply integrated with VS Code.",
                "support_score": 0.78,
                "accepted": True,
                "evidence_ids": ["E-002"],
            },
            {
                "claim_id": "C-003",
                "competitor_name": "Tabnine",
                "dimension": "positioning",
                "statement": "Tabnine emphasizes enterprise privacy and deployment control.",
                "support_score": 0.65,
                "accepted": True,
                "evidence_ids": ["E-003"],
            },
        ],
    )
    _write_json(
        tmp_path / "evidence.json",
        [
            {
                "evidence_id": "E-001",
                "source_id": "cursor.com/pricing",
                "competitor_name": "Cursor",
                "quote": "Pro plan",
            },
            {
                "evidence_id": "E-002",
                "source_id": "github.com/features/copilot",
                "competitor_name": "Copilot",
                "quote": "Works in VS Code",
            },
        ],
    )
    _write_json(tmp_path / "source_summary.json", {"total": 6, "fetched": 4, "valid": 4})
    return tmp_path


@pytest.fixture()
def fake_png_renderer(monkeypatch):
    def render(slide_defs: list[tuple[str, str]], run_dir: Path) -> list[Path]:
        slides_dir = Path(run_dir) / "_slides"
        slides_dir.mkdir(parents=True, exist_ok=True)
        paths: list[Path] = []
        for index, (template_name, html) in enumerate(slide_defs, start=1):
            stem = Path(template_name).stem
            (slides_dir / f"{index:02d}_{stem}.html").write_text(
                html,
                encoding="utf-8",
            )
            png_path = slides_dir / f"{index:02d}_{stem}.png"
            png_path.write_bytes(_PNG_BYTES)
            paths.append(png_path)
        return paths

    monkeypatch.setattr(ppt_visual_builder, "_render_slides_to_png", render)
    return render


def test_build_ppt_visual_creates_file(
    minimal_run_dir: Path,
    fake_png_renderer,
) -> None:
    path = build_ppt_visual(minimal_run_dir)

    assert path == minimal_run_dir / "report.pptx"
    assert path.exists()
    assert path.stat().st_size > 5000


def test_ppt_visual_has_five_slides(
    minimal_run_dir: Path,
    fake_png_renderer,
) -> None:
    path = build_ppt_visual(minimal_run_dir)

    prs = Presentation(path)

    assert len(prs.slides) == 5


def test_each_slide_has_full_page_picture(
    minimal_run_dir: Path,
    fake_png_renderer,
) -> None:
    path = build_ppt_visual(minimal_run_dir)

    prs = Presentation(path)

    for slide in prs.slides:
        pictures = [
            shape
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        assert len(pictures) == 1
        picture = pictures[0]
        assert picture.left == 0
        assert picture.top == 0
        assert picture.width == prs.slide_width
        assert picture.height == prs.slide_height


def test_visual_builder_writes_slide_workspace(
    minimal_run_dir: Path,
    fake_png_renderer,
) -> None:
    build_ppt_visual(minimal_run_dir)

    slides_dir = minimal_run_dir / "_slides"

    assert len(list(slides_dir.glob("*.html"))) == 5
    assert len(list(slides_dir.glob("*.png"))) == 5


def test_finalize_falls_back_to_text_ppt_on_visual_failure(
    tmp_path: Path,
    monkeypatch,
) -> None:
    import cs_mvp.agents.ppt_builder as text_builder
    import cs_mvp.graph as graph_module

    db_path = str(tmp_path / "cs_mvp.db")
    runs_dir = tmp_path / "runs"
    db.init_db(db_path)

    task = AnalysisTask(
        task_id="T-ppt-fallback",
        query="AI IDE",
        competitors=[CompetitorInput(name="Cursor")],
    )
    run = AgentRun(run_id="RUN-ppt-fallback", task_id=task.task_id)
    db.insert_task(task)
    db.insert_run(run)
    monkeypatch.setattr(graph_module, "_RUNS_DIR", str(runs_dir))

    def fail_visual(run_dir: Path) -> Path:
        raise RuntimeError("playwright unavailable")

    calls: list[Path] = []

    def build_text(run_dir: Path) -> Path:
        calls.append(Path(run_dir))
        pptx_path = Path(run_dir) / "report.pptx"
        pptx_path.write_bytes(b"text ppt fallback")
        return pptx_path

    monkeypatch.setattr(ppt_visual_builder, "build_ppt_visual", fail_visual)
    monkeypatch.setattr(text_builder, "build_ppt", build_text)

    graph_module.node_finalize(
        GraphState(
            task=task,
            run_id=run.run_id,
            report_md="# Report\n\nFallback smoke.",
        )
    )

    run_dir = runs_dir / task.task_id
    assert calls == [run_dir]
    assert (run_dir / "report.pptx").read_bytes() == b"text ppt fallback"
