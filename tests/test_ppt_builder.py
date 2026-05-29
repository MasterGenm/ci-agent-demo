from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation

from cs_mvp.agents.ppt_builder import build_ppt


def _write_json(path: Path, payload: object) -> None:
    path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")


@pytest.fixture()
def minimal_run_dir(tmp_path: Path) -> Path:
    run_summary = {
        "task_id": "T-test-001",
        "run_id": "RUN-test-001",
        "query": "AI 代码编辑器竞品分析",
        "competitors": ["Cursor", "Copilot", "Tabnine"],
        "status": "completed",
        "completed_at": "2026-05-25 12:00:00",
        "duration_seconds": 300.0,
    }
    claims = [
        {
            "claim_id": "C-001",
            "competitor_name": "Cursor",
            "dimension": "pricing",
            "statement": "Cursor 提供 $20/月的 Pro 计划，含 500 次快速请求。",
            "support_score": 0.85,
            "accepted": True,
            "evidence_ids": ["E-001"],
        },
        {
            "claim_id": "C-002",
            "competitor_name": "Copilot",
            "dimension": "features",
            "statement": "GitHub Copilot 深度集成 VS Code，支持多文件上下文。",
            "support_score": 0.78,
            "accepted": True,
            "evidence_ids": ["E-002"],
        },
        {
            "claim_id": "C-003",
            "competitor_name": "Tabnine",
            "dimension": "positioning",
            "statement": "Tabnine 更强调企业私有化和代码隐私。",
            "support_score": 0.65,
            "accepted": True,
            "evidence_ids": ["E-003"],
        },
    ]
    evidence = [
        {
            "evidence_id": "E-001",
            "source_id": "cursor.com/pricing",
            "competitor_name": "Cursor",
            "quote": "Pro plan at $20/month",
        },
        {
            "evidence_id": "E-002",
            "source_id": "github.com/features/copilot",
            "competitor_name": "Copilot",
            "quote": "Copilot works in VS Code.",
        },
    ]
    source_summary = {"total": 6, "fetched": 4, "valid": 4}
    _write_json(tmp_path / "run_summary.json", run_summary)
    _write_json(tmp_path / "claims.json", claims)
    _write_json(tmp_path / "evidence.json", evidence)
    _write_json(tmp_path / "source_summary.json", source_summary)
    return tmp_path


def _slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            parts.append(shape.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "\n".join(parts)


def test_build_ppt_creates_file(minimal_run_dir: Path) -> None:
    path = build_ppt(minimal_run_dir)

    assert path == minimal_run_dir / "report.pptx"
    assert path.exists()
    assert path.stat().st_size > 0


def test_ppt_has_five_slides(minimal_run_dir: Path) -> None:
    path = build_ppt(minimal_run_dir)

    prs = Presentation(path)

    assert len(prs.slides) == 5


def test_ppt_cover_contains_query(minimal_run_dir: Path) -> None:
    path = build_ppt(minimal_run_dir)

    prs = Presentation(path)
    text = _slide_text(prs.slides[0])

    assert "AI 代码编辑器竞品分析" in text


def test_ppt_matrix_slide_has_competitor_names(minimal_run_dir: Path) -> None:
    path = build_ppt(minimal_run_dir)

    prs = Presentation(path)
    text = _slide_text(prs.slides[2])

    assert "Cursor" in text
    assert "Copilot" in text
    assert "Tabnine" in text


def test_build_ppt_handles_empty_claims(minimal_run_dir: Path) -> None:
    _write_json(minimal_run_dir / "claims.json", [])

    path = build_ppt(minimal_run_dir)
    prs = Presentation(path)

    assert path.exists()
    assert len(prs.slides) == 5
